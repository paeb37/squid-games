#!/usr/bin/env python3
"""
PowerPoint Information Extractor using python-pptx
This script extracts various types of information from PowerPoint files.
"""

from pptx import Presentation
import os
import json
from typing import Dict, List, Any
import argparse


class PowerPointExtractor:
    """Extract information from PowerPoint files using python-pptx library."""
    
    def __init__(self, file_path: str):
        """Initialize the extractor with a PowerPoint file."""
        self.file_path = file_path
        self.presentation = None
        self.load_presentation()
    
    def load_presentation(self):
        """Load the PowerPoint presentation."""
        try:
            self.presentation = Presentation(self.file_path)
            print(f"Successfully loaded PowerPoint file: {self.file_path}")
        except Exception as e:
            raise Exception(f"Error loading PowerPoint file: {e}")
    
    def extract_basic_info(self) -> Dict[str, Any]:
        """Extract basic information about the presentation."""
        info = {
            "file_path": self.file_path,
            "total_slides": len(self.presentation.slides),
            "slide_dimensions": {
                "width": self.presentation.slide_width,
                "height": self.presentation.slide_height
            }
        }
        return info
    
    def extract_all_text(self) -> Dict[str, List[str]]:
        """Extract all text from all slides."""
        all_text = {
            "slides": [],
            "all_text_combined": []
        }
        
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slide_info = {
                "slide_number": i,
                "text_elements": slide_text,
                "combined_text": " ".join(slide_text)
            }
            
            all_text["slides"].append(slide_info)
            all_text["all_text_combined"].extend(slide_text)
        
        return all_text
    
    def extract_slide_titles(self) -> List[Dict[str, Any]]:
        """Extract titles from all slides."""
        titles = []
        
        for i, slide in enumerate(self.presentation.slides, 1):
            title = "No Title"
            
            # Try to find the title placeholder
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # Title placeholder
                        if hasattr(shape, "text") and shape.text.strip():
                            title = shape.text.strip()
                        break
            
            # If no title placeholder found, use the first text shape
            if title == "No Title":
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip().split('\n')[0]  # Take first line
                        break
            
            titles.append({
                "slide_number": i,
                "title": title
            })
        
        return titles
    
    def extract_images_info(self) -> List[Dict[str, Any]]:
        """Extract information about images in the presentation."""
        images_info = []
        
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_images = []
            
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture shape type
                    image_info = {
                        "shape_name": shape.name,
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top
                    }
                    slide_images.append(image_info)
            
            if slide_images:
                images_info.append({
                    "slide_number": i,
                    "images": slide_images,
                    "image_count": len(slide_images)
                })
        
        return images_info
    
    def extract_notes(self) -> List[Dict[str, Any]]:
        """Extract speaker notes from slides."""
        notes_info = []
        
        for i, slide in enumerate(self.presentation.slides, 1):
            notes_slide = slide.notes_slide
            notes_text = ""
            
            if notes_slide:
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        notes_text += shape.text.strip() + " "
            
            if notes_text.strip():
                notes_info.append({
                    "slide_number": i,
                    "notes": notes_text.strip()
                })
        
        return notes_info
    
    def extract_slide_layout_info(self) -> List[Dict[str, Any]]:
        """Extract layout information for each slide."""
        layout_info = []
        
        for i, slide in enumerate(self.presentation.slides, 1):
            layout = slide.slide_layout
            
            slide_layout_info = {
                "slide_number": i,
                "layout_name": layout.name,
                "shape_count": len(slide.shapes),
                "shapes": []
            }
            
            # Get information about each shape
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "has_text": hasattr(shape, "text") and bool(shape.text.strip()),
                    "is_placeholder": shape.is_placeholder
                }
                slide_layout_info["shapes"].append(shape_info)
            
            layout_info.append(slide_layout_info)
        
        return layout_info
    
    def extract_all_information(self) -> Dict[str, Any]:
        """Extract all available information from the PowerPoint file."""
        return {
            "basic_info": self.extract_basic_info(),
            "titles": self.extract_slide_titles(),
            "text_content": self.extract_all_text(),
            "images_info": self.extract_images_info(),
            "notes": self.extract_notes(),
            "layout_info": self.extract_slide_layout_info()
        }
    
    def save_to_json(self, output_file: str = None):
        """Save extracted information to a JSON file."""
        if not output_file:
            base_name = os.path.splitext(os.path.basename(self.file_path))[0]
            output_file = os.path.join('slides', f"{base_name}_extracted_info.json")
        
        all_info = self.extract_all_information()
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(all_info, f, indent=2, ensure_ascii=False)
        
        print(f"Extracted information saved to: {output_file}")
        return output_file


def main():
    """Main function to run the PowerPoint extractor."""
    parser = argparse.ArgumentParser(description="Extract information from PowerPoint files")
    parser.add_argument("file_path", help="Path to the PowerPoint file (.pptx)")
    parser.add_argument("--output", "-o", help="Output JSON file path", default=None)
    parser.add_argument("--text-only", action="store_true", help="Extract only text content")
    parser.add_argument("--titles-only", action="store_true", help="Extract only slide titles")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.file_path):
        print(f"Error: File '{args.file_path}' does not exist.")
        return
    
    try:
        extractor = PowerPointExtractor(args.file_path)
        
        if args.text_only:
            text_info = extractor.extract_all_text()
            print("\n=== EXTRACTED TEXT ===")
            for slide in text_info["slides"]:
                print(f"\nSlide {slide['slide_number']}:")
                print(slide['combined_text'])
        
        elif args.titles_only:
            titles = extractor.extract_slide_titles()
            print("\n=== SLIDE TITLES ===")
            for title_info in titles:
                print(f"Slide {title_info['slide_number']}: {title_info['title']}")
        
        else:
            # Extract all information and save to JSON
            output_file = extractor.save_to_json(args.output)
            
            # Also print a summary
            basic_info = extractor.extract_basic_info()
            print(f"\n=== SUMMARY ===")
            print(f"File: {basic_info['file_path']}")
            print(f"Total slides: {basic_info['total_slides']}")
            print(f"Dimensions: {basic_info['slide_dimensions']['width']} x {basic_info['slide_dimensions']['height']}")
            print(f"Full extraction saved to: {output_file}")
    
    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    main()
